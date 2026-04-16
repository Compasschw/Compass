"""Build Compass x Pear Suite meeting deck as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colours ──────────────────────────────────────────────
COMPASS_GREEN = RGBColor(0x0D, 0x9B, 0x6A)
COMPASS_DARK = RGBColor(0x06, 0x4E, 0x3B)
PEAR_GOLD = RGBColor(0xF5, 0x9E, 0x0B)
PEAR_DARK = RGBColor(0x92, 0x40, 0x0E)
SLATE_900 = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700 = RGBColor(0x33, 0x41, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_300 = RGBColor(0xCB, 0xD5, 0xE1)
SLATE_100 = RGBColor(0xF1, 0xF5, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_500 = RGBColor(0xEF, 0x44, 0x44)
BLUE_500 = RGBColor(0x3B, 0x82, 0xF6)
BLUE_DARK = RGBColor(0x1D, 0x4E, 0xD8)
PURPLE_500 = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_DARK = RGBColor(0x6D, 0x28, 0xD9)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
GOLD_BG = RGBColor(0xFF, 0xF7, 0xED)
BLUE_BG = RGBColor(0xEF, 0xF6, 0xFF)
PURPLE_BG = RGBColor(0xF5, 0xF3, 0xFF)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────
def add_shape(slide, left, top, width, height, fill=None, border_color=None,
              border_width=Pt(1.5), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.shadow.inherit = False
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if border_color:
        shp.line.color.rgb = border_color
        shp.line.width = border_width
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, left, top, width, height, text, font_size=12,
             bold=False, color=SLATE_700, alignment=PP_ALIGN.LEFT,
             font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_circle(slide, left, top, diameter, fill_color, text="",
               font_size=12, text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    if text:
        tf = shp.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shp


def add_arrow(slide, start_left, start_top, end_left, end_top,
              color=SLATE_500, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1, start_left, start_top, end_left, end_top  # 1 = straight
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def add_right_arrow_shape(slide, left, top, width, height, fill_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp


def add_chevron(slide, left, top, width, height, fill_color):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    shp.line.fill.background()
    return shp


def make_dark_slide(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = SLATE_900


def make_white_slide(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = WHITE


def add_q_badge(slide, number, left, top):
    add_circle(slide, left, top, Inches(0.45), COMPASS_GREEN, str(number), 14, WHITE)


def add_category(slide, text, left, top):
    add_text(slide, left, top, Inches(3), Inches(0.3), text.upper(),
             font_size=10, bold=True, color=COMPASS_GREEN)


def add_question_title(slide, text, left, top, width=Inches(10)):
    add_text(slide, left, top, width, Inches(0.8), text,
             font_size=26, bold=True, color=SLATE_900)


def add_context_box(slide, left, top, width, height, label, body_lines,
                    bg=SLATE_100, label_color=SLATE_500, text_color=SLATE_700):
    box = add_shape(slide, left, top, width, height, fill=bg, border_color=None)
    add_text(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4),
             Inches(0.22), label.upper(), font_size=9, bold=True, color=label_color)
    y = top + Inches(0.4)
    for line in body_lines:
        bullet = "• " if not line.startswith("•") else ""
        add_text(slide, left + Inches(0.2), y, width - Inches(0.4), Inches(0.22),
                 f"{bullet}{line}", font_size=11, color=text_color)
        y += Inches(0.24)
    return box


def add_highlight_box(slide, left, top, width, height, label, body_text,
                      bg=GREEN_BG, border=COMPASS_GREEN, label_color=COMPASS_GREEN,
                      text_color=COMPASS_DARK):
    box = add_shape(slide, left, top, width, height, fill=bg, border_color=border, border_width=Pt(3))
    add_text(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4),
             Inches(0.22), label.upper(), font_size=9, bold=True, color=label_color)
    add_text(slide, left + Inches(0.2), top + Inches(0.4), width - Inches(0.4),
             height - Inches(0.5), body_text, font_size=11, color=text_color)
    return box


# ============================================================
# SLIDE 0: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
make_dark_slide(slide)

# Compass logo badge
add_shape(slide, Inches(5.15), Inches(1.8), Inches(0.9), Inches(0.9),
          fill=COMPASS_GREEN, border_color=None)
add_text(slide, Inches(5.15), Inches(1.88), Inches(0.9), Inches(0.8), "C",
         font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# × symbol
add_text(slide, Inches(6.2), Inches(1.88), Inches(0.6), Inches(0.8), "×",
         font_size=28, color=SLATE_500, alignment=PP_ALIGN.CENTER)

# Pear Suite logo badge
add_shape(slide, Inches(7.0), Inches(1.8), Inches(0.9), Inches(0.9),
          fill=PEAR_GOLD, border_color=None)
add_text(slide, Inches(7.0), Inches(1.88), Inches(0.9), Inches(0.8), "P",
         font_size=36, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)

# Title
add_text(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.8),
         "Compass × Pear Suite", font_size=44, bold=True, color=WHITE,
         alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(4.0), Inches(10), Inches(0.5),
         "Integration Discovery — Technical Discussion", font_size=20,
         color=SLATE_300, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.4),
         "April 2026", font_size=14, color=SLATE_500, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 1: Vision & Alignment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 1, Inches(0.6), Inches(0.4))
add_category(slide, "Strategy & Alignment", Inches(1.15), Inches(0.45))
add_question_title(slide, "What's Pear Suite's vision for the CHW ecosystem,\nand where do you see the biggest gaps today?",
                   Inches(0.6), Inches(1.0))

# Context card
add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.5),
                "Why This Matters", [
                    "Compass fills the session documentation",
                    "& field workflow gap for CHWs.",
                    "Need to understand where Pear Suite",
                    "sees unmet needs."
                ])

add_highlight_box(slide, Inches(0.6), Inches(4.0), Inches(3.5), Inches(1.2),
                  "Our Thesis",
                  "Pear Suite owns billing infrastructure.\nCompass owns the CHW experience.\nTogether: field visit → paid claim.")

# Venn diagram area
# Compass circle
add_shape(slide, Inches(5.0), Inches(2.2), Inches(3.8), Inches(4.0),
          fill=GREEN_BG, border_color=COMPASS_GREEN, shape_type=MSO_SHAPE.OVAL)
# Pear Suite circle
add_shape(slide, Inches(7.2), Inches(2.2), Inches(3.8), Inches(4.0),
          fill=GOLD_BG, border_color=PEAR_GOLD, shape_type=MSO_SHAPE.OVAL)

# Labels
add_text(slide, Inches(5.3), Inches(2.6), Inches(2), Inches(0.3), "Compass",
         font_size=14, bold=True, color=COMPASS_DARK)
add_text(slide, Inches(5.3), Inches(3.0), Inches(2), Inches(0.25), "• Session docs", font_size=11, color=SLATE_700)
add_text(slide, Inches(5.3), Inches(3.3), Inches(2), Inches(0.25), "• CHW workflows", font_size=11, color=SLATE_700)
add_text(slide, Inches(5.3), Inches(3.6), Inches(2), Inches(0.25), "• Mobile-first UX", font_size=11, color=SLATE_700)
add_text(slide, Inches(5.3), Inches(3.9), Inches(2), Inches(0.25), "• Field data capture", font_size=11, color=SLATE_700)

add_text(slide, Inches(9.0), Inches(2.6), Inches(2), Inches(0.3), "Pear Suite",
         font_size=14, bold=True, color=PEAR_DARK)
add_text(slide, Inches(9.0), Inches(3.0), Inches(2), Inches(0.25), "• Claims submission", font_size=11, color=SLATE_700)
add_text(slide, Inches(9.0), Inches(3.3), Inches(2), Inches(0.25), "• Clearinghouse", font_size=11, color=SLATE_700)
add_text(slide, Inches(9.0), Inches(3.6), Inches(2), Inches(0.25), "• EHR integration", font_size=11, color=SLATE_700)
add_text(slide, Inches(9.0), Inches(3.9), Inches(2), Inches(0.25), "• Medi-Cal billing", font_size=11, color=SLATE_700)

# Overlap
add_text(slide, Inches(7.4), Inches(3.5), Inches(2), Inches(0.3), "Overlap",
         font_size=13, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(7.4), Inches(3.8), Inches(2), Inches(0.25), "CHW encounter",
         font_size=11, color=SLATE_700, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(7.4), Inches(4.05), Inches(2), Inches(0.25), "data exchange",
         font_size=11, color=SLATE_700, alignment=PP_ALIGN.CENTER)

# Bottom tag
add_shape(slide, Inches(6.5), Inches(6.5), Inches(3), Inches(0.4),
          fill=COMPASS_GREEN, border_color=None)
add_text(slide, Inches(6.5), Inches(6.5), Inches(3), Inches(0.4),
         "Field visit → Paid claim", font_size=13, bold=True, color=WHITE,
         alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: Compass Session Flow — One-Pager
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_circle(slide, Inches(0.6), Inches(0.4), Inches(0.45), SLATE_900, "⚙", 14, WHITE)
add_category(slide, "How Compass Works Today", Inches(1.15), Inches(0.45))
add_question_title(slide, "Compass Session Documentation Flow — and Where Pear Suite Fits",
                   Inches(0.6), Inches(1.0), Inches(11))

# ── Compass zone ──
add_shape(slide, Inches(0.4), Inches(1.8), Inches(7.6), Inches(5.4),
          fill=GREEN_BG, border_color=COMPASS_GREEN, border_width=Pt(2))
add_text(slide, Inches(0.6), Inches(1.9), Inches(2), Inches(0.3),
         "COMPASS PLATFORM", font_size=12, bold=True, color=COMPASS_DARK)

# Steps row 1 (top)
step_w = Inches(2.2)
step_h = Inches(1.0)
row1_y = Inches(2.4)
steps_top = [
    ("1", "Start Session", "CHW opens app\nSelects or enrolls member"),
    ("2", "Member Lookup", "Verify Medi-Cal ID\nConfirm eligibility & plan"),
    ("3", "Service Details", "Session type & location\nDate, start time, duration"),
]
for i, (num, title, desc) in enumerate(steps_top):
    x = Inches(0.6) + i * Inches(2.5)
    add_shape(slide, x, row1_y, step_w, step_h, fill=WHITE, border_color=COMPASS_GREEN)
    add_circle(slide, x + Inches(0.1), row1_y + Inches(0.1), Inches(0.3), COMPASS_GREEN, num, 11)
    add_text(slide, x + Inches(0.5), row1_y + Inches(0.08), Inches(1.6), Inches(0.25),
             title, font_size=12, bold=True, color=COMPASS_DARK)
    add_text(slide, x + Inches(0.15), row1_y + Inches(0.42), step_w - Inches(0.3), Inches(0.55),
             desc, font_size=10, color=SLATE_700)
    # Arrow between
    if i < 2:
        add_right_arrow_shape(slide, x + step_w + Inches(0.05), row1_y + Inches(0.35),
                              Inches(0.2), Inches(0.25), COMPASS_GREEN)

# Steps row 2 (bottom, reversed direction)
row2_y = Inches(3.7)
steps_bot = [
    ("6", "Data Package", "Member ID + demographics\nService type + duration\nSession notes + referrals"),
    ("5", "Complete Session", "CHW reviews & signs off\nAuto-generates structured\nsession documentation"),
    ("4", "Document Visit", "Session notes / narrative\nGoals addressed\nReferrals & follow-ups"),
]
for i, (num, title, desc) in enumerate(steps_bot):
    x = Inches(0.6) + i * Inches(2.5)
    bc = COMPASS_GREEN
    bw = Pt(2.5) if num == "6" else Pt(1.5)
    add_shape(slide, x, row2_y, step_w, Inches(1.15), fill=WHITE, border_color=bc, border_width=bw)
    add_circle(slide, x + Inches(0.1), row2_y + Inches(0.1), Inches(0.3), COMPASS_GREEN, num, 11)
    add_text(slide, x + Inches(0.5), row2_y + Inches(0.08), Inches(1.6), Inches(0.25),
             title, font_size=12, bold=True, color=COMPASS_DARK)
    add_text(slide, x + Inches(0.15), row2_y + Inches(0.42), step_w - Inches(0.3), Inches(0.7),
             desc, font_size=10, color=SLATE_700)
    if i < 2:
        arr_x = x + step_w + Inches(0.05)
        # Reverse arrows (pointing left)
        add_shape(slide, arr_x, row2_y + Inches(0.4), Inches(0.22), Inches(0.25),
                  fill=COMPASS_GREEN, border_color=None,
                  shape_type=MSO_SHAPE.LEFT_ARROW)

# Data fields table
tbl_y = Inches(5.1)
add_shape(slide, Inches(0.6), tbl_y, Inches(7.2), Inches(1.9),
          fill=WHITE, border_color=COMPASS_GREEN, border_width=Pt(1.5))
add_text(slide, Inches(0.8), tbl_y + Inches(0.08), Inches(3), Inches(0.25),
         "Data Captured Per Session", font_size=12, bold=True, color=COMPASS_DARK)

cols = [
    ("MEMBER", ["Medi-Cal ID", "Full name & DOB", "Address / county", "Managed care plan", "Contact info"]),
    ("SERVICE", ["Session type", "Date & duration", "Place of service", "CHW provider info", "Org / NPI (config)"]),
    ("DOCUMENTATION", ["Visit narrative", "Goals addressed", "Referrals made", "Follow-up actions", "Supervisor sign-off"]),
]
for ci, (label, items) in enumerate(cols):
    cx = Inches(0.8) + ci * Inches(2.4)
    add_text(slide, cx, tbl_y + Inches(0.35), Inches(2), Inches(0.2),
             label, font_size=9, bold=True, color=SLATE_500)
    for ri, item in enumerate(items):
        add_text(slide, cx, tbl_y + Inches(0.55) + ri * Inches(0.22), Inches(2.2), Inches(0.2),
                 f"• {item}", font_size=10, color=SLATE_700)

# ── Integration seam ──
add_shape(slide, Inches(8.25), Inches(1.8), Inches(0.25), Inches(5.4),
          fill=SLATE_900, border_color=None, shape_type=MSO_SHAPE.RECTANGLE)
txb = add_text(slide, Inches(8.05), Inches(4.0), Inches(0.7), Inches(1.5),
               "INTEGRATION\nSEAM", font_size=9, bold=True, color=WHITE,
               alignment=PP_ALIGN.CENTER)

# Big arrow from Compass → seam
add_right_arrow_shape(slide, Inches(7.85), Inches(3.3), Inches(0.4), Inches(0.3), COMPASS_GREEN)

# ── Pear Suite zone ──
add_shape(slide, Inches(8.8), Inches(1.8), Inches(4.2), Inches(5.4),
          fill=GOLD_BG, border_color=PEAR_GOLD, border_width=Pt(2))
add_text(slide, Inches(9.0), Inches(1.9), Inches(2), Inches(0.3),
         "PEAR SUITE", font_size=12, bold=True, color=PEAR_DARK)
add_text(slide, Inches(9.0), Inches(2.15), Inches(3), Inches(0.25),
         "Billing & claims infrastructure", font_size=10, color=SLATE_500)

# Arrow from seam → Pear Suite
add_right_arrow_shape(slide, Inches(8.55), Inches(3.3), Inches(0.4), Inches(0.3), PEAR_GOLD)

ps_steps = [
    ("A", "Receive Session Data", "API ingests Compass data package"),
    ("B", "Map & Validate", "Add CPT/HCPCS, ICD-10, apply rules"),
    ("C", "Submit to Clearinghouse", "Route claim to Medi-Cal / plan"),
    ("D", "Adjudicate & Report", "Paid / denied / resubmit status"),
]
for i, (letter, title, desc) in enumerate(ps_steps):
    y = Inches(2.6) + i * Inches(1.15)
    add_shape(slide, Inches(9.0), y, Inches(3.8), Inches(0.85),
              fill=WHITE, border_color=PEAR_GOLD)
    add_circle(slide, Inches(9.15), y + Inches(0.15), Inches(0.3), PEAR_GOLD, letter, 11, WHITE)
    add_text(slide, Inches(9.55), y + Inches(0.1), Inches(3), Inches(0.25),
             title, font_size=12, bold=True, color=PEAR_DARK)
    add_text(slide, Inches(9.55), y + Inches(0.4), Inches(3), Inches(0.25),
             desc, font_size=10, color=SLATE_700)
    if i < 3:
        add_shape(slide, Inches(10.7), y + Inches(0.85), Inches(0.25), Inches(0.2),
                  fill=PEAR_GOLD, border_color=None,
                  shape_type=MSO_SHAPE.DOWN_ARROW)

# Status callback label
add_text(slide, Inches(8.8), Inches(7.0), Inches(4), Inches(0.25),
         "↩ Status callback to Compass (webhook / polling)", font_size=10,
         bold=True, color=PURPLE_500, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: Current Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 2, Inches(0.6), Inches(0.4))
add_category(slide, "Current State", Inches(1.15), Inches(0.45))
add_question_title(slide, "How are CHW organizations currently getting session data\ninto Pear Suite for claims submission?",
                   Inches(0.6), Inches(1.0))

# Context
add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.3),
                "Why We're Asking", [
                    "Need to understand existing data flow",
                    "before proposing a new one.",
                    "Manual entry? CSV? API?"
                ])

add_highlight_box(slide, Inches(0.6), Inches(3.8), Inches(3.5), Inches(1.1),
                  "Assumption to Validate",
                  "Most CHW orgs rely on manual data entry\nor paper forms → significant time gap.",
                  bg=GOLD_BG, border=PEAR_GOLD, label_color=PEAR_DARK, text_color=PEAR_DARK)

# Current flow (top)
diagram_x = Inches(4.8)
add_text(slide, diagram_x, Inches(2.2), Inches(6), Inches(0.35),
         "Current Suspected Workflow", font_size=15, bold=True, color=SLATE_900)

flow_y = Inches(2.7)
boxes_cur = [
    ("CHW", "Field visit", GREEN_BG, COMPASS_GREEN),
    ("Paper / Notes", "Manual capture", GOLD_BG, PEAR_GOLD),
    ("Manual Entry", "Into Pear Suite", RED_BG, RED_500),
]
for i, (title, sub, bg, bc) in enumerate(boxes_cur):
    x = diagram_x + i * Inches(2.6)
    add_shape(slide, x, flow_y, Inches(2.1), Inches(0.85), fill=bg, border_color=bc)
    add_text(slide, x + Inches(0.15), flow_y + Inches(0.1), Inches(1.8), Inches(0.25),
             title, font_size=12, bold=True, color=SLATE_900)
    add_text(slide, x + Inches(0.15), flow_y + Inches(0.4), Inches(1.8), Inches(0.25),
             sub, font_size=11, color=SLATE_700)
    if i < 2:
        add_right_arrow_shape(slide, x + Inches(2.15), flow_y + Inches(0.25),
                              Inches(0.35), Inches(0.3), SLATE_500)

add_text(slide, diagram_x + Inches(2.2), flow_y + Inches(0.95), Inches(1), Inches(0.2),
         "hours/days", font_size=9, bold=True, color=RED_500, alignment=PP_ALIGN.CENTER)
add_text(slide, diagram_x + Inches(4.8), flow_y + Inches(0.95), Inches(1), Inches(0.2),
         "hours/days", font_size=9, bold=True, color=RED_500, alignment=PP_ALIGN.CENTER)

# Divider
add_shape(slide, diagram_x, Inches(4.1), Inches(7.8), Inches(0.02),
          fill=SLATE_300, border_color=None, shape_type=MSO_SHAPE.RECTANGLE)
add_text(slide, diagram_x, Inches(4.2), Inches(6), Inches(0.35),
         "Proposed: Compass Integration", font_size=15, bold=True, color=COMPASS_GREEN)

# Proposed flow
flow2_y = Inches(4.7)
boxes_new = [
    ("CHW", "Field visit", GREEN_BG, COMPASS_GREEN),
    ("Compass App", "Digital capture\nin real-time", GREEN_BG, COMPASS_GREEN),
    ("Pear Suite API", "Auto-submit\nclaims", GOLD_BG, PEAR_GOLD),
]
for i, (title, sub, bg, bc) in enumerate(boxes_new):
    x = diagram_x + i * Inches(2.6)
    add_shape(slide, x, flow2_y, Inches(2.1), Inches(0.85), fill=bg, border_color=bc, border_width=Pt(2.5))
    add_text(slide, x + Inches(0.15), flow2_y + Inches(0.1), Inches(1.8), Inches(0.25),
             title, font_size=12, bold=True, color=SLATE_900)
    add_text(slide, x + Inches(0.15), flow2_y + Inches(0.4), Inches(1.8), Inches(0.4),
             sub, font_size=11, color=SLATE_700)
    if i < 2:
        add_right_arrow_shape(slide, x + Inches(2.15), flow2_y + Inches(0.25),
                              Inches(0.35), Inches(0.3), COMPASS_GREEN)

add_text(slide, diagram_x + Inches(2.2), flow2_y + Inches(0.95), Inches(1), Inches(0.2),
         "real-time", font_size=9, bold=True, color=COMPASS_GREEN, alignment=PP_ALIGN.CENTER)
add_text(slide, diagram_x + Inches(4.8), flow2_y + Inches(0.95), Inches(1), Inches(0.2),
         "minutes", font_size=9, bold=True, color=COMPASS_GREEN, alignment=PP_ALIGN.CENTER)

# Bottom pill
add_shape(slide, Inches(7.0), Inches(6.3), Inches(3), Inches(0.45),
          fill=COMPASS_GREEN, border_color=None)
add_text(slide, Inches(7.0), Inches(6.33), Inches(3), Inches(0.4),
         "Days → Minutes", font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4: API Surface
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 3, Inches(0.6), Inches(0.4))
add_category(slide, "Technical Architecture", Inches(1.15), Inches(0.45))
add_question_title(slide, "What does your API surface look like — REST, GraphQL,\nwebhooks, or file-based exchange?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.3),
                "Why This Matters", [
                    "Determines entire integration",
                    "architecture and dev timeline.",
                    "Each approach has very different",
                    "real-time vs. batch implications."
                ])

add_context_box(slide, Inches(0.6), Inches(3.8), Inches(3.5), Inches(1.3),
                "We Can Support", [
                    "REST API (preferred)",
                    "Webhook callbacks",
                    "GraphQL",
                    "Batch file exchange (SFTP/S3)"
                ])

# Diagram: two boxes with 4 connection options
# Compass box
add_shape(slide, Inches(4.8), Inches(2.5), Inches(2.2), Inches(3.2),
          fill=GREEN_BG, border_color=COMPASS_GREEN, border_width=Pt(2))
add_text(slide, Inches(4.8), Inches(2.6), Inches(2.2), Inches(0.3),
         "Compass", font_size=14, bold=True, color=COMPASS_DARK, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(5.0), Inches(3.0), Inches(1.8), Inches(0.22), "Session data", font_size=11, color=SLATE_700)
add_text(slide, Inches(5.0), Inches(3.3), Inches(1.8), Inches(0.22), "CHW encounters", font_size=11, color=SLATE_700)
add_text(slide, Inches(5.0), Inches(3.6), Inches(1.8), Inches(0.22), "Member info", font_size=11, color=SLATE_700)
add_text(slide, Inches(5.0), Inches(3.9), Inches(1.8), Inches(0.22), "Service codes", font_size=11, color=SLATE_700)
add_text(slide, Inches(4.8), Inches(5.1), Inches(2.2), Inches(0.22),
         "Python / FastAPI", font_size=10, color=SLATE_500, alignment=PP_ALIGN.CENTER)

# Pear Suite box
add_shape(slide, Inches(10.0), Inches(2.5), Inches(2.2), Inches(3.2),
          fill=GOLD_BG, border_color=PEAR_GOLD, border_width=Pt(2))
add_text(slide, Inches(10.0), Inches(2.6), Inches(2.2), Inches(0.3),
         "Pear Suite", font_size=14, bold=True, color=PEAR_DARK, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(10.2), Inches(3.0), Inches(1.8), Inches(0.22), "Claims engine", font_size=11, color=SLATE_700)
add_text(slide, Inches(10.2), Inches(3.3), Inches(1.8), Inches(0.22), "Clearinghouse", font_size=11, color=SLATE_700)
add_text(slide, Inches(10.2), Inches(3.6), Inches(1.8), Inches(0.22), "EHR connectors", font_size=11, color=SLATE_700)
add_text(slide, Inches(10.2), Inches(3.9), Inches(1.8), Inches(0.22), "Medi-Cal rules", font_size=11, color=SLATE_700)

# Connection options
options = [
    ("A: REST API", BLUE_BG, BLUE_500, BLUE_DARK),
    ("B: Webhooks", PURPLE_BG, PURPLE_500, PURPLE_DARK),
    ("C: GraphQL", GREEN_BG, COMPASS_GREEN, COMPASS_DARK),
    ("D: Batch / SFTP", GOLD_BG, PEAR_GOLD, PEAR_DARK),
]
for i, (label, bg, bc, tc) in enumerate(options):
    y = Inches(2.7) + i * Inches(0.7)
    add_shape(slide, Inches(7.3), y, Inches(2.4), Inches(0.5), fill=bg, border_color=bc)
    add_text(slide, Inches(7.3), y + Inches(0.05), Inches(2.4), Inches(0.35),
             label, font_size=12, bold=True, color=tc, alignment=PP_ALIGN.CENTER)
    # Left arrow
    add_right_arrow_shape(slide, Inches(7.0), y + Inches(0.12), Inches(0.25), Inches(0.22), bc)
    # Right arrow
    add_right_arrow_shape(slide, Inches(9.75), y + Inches(0.12), Inches(0.25), Inches(0.22), bc)

add_text(slide, Inches(7.3), Inches(5.8), Inches(2.4), Inches(0.4),
         "Preferred: real-time API (A or C)\nFallback: webhook + batch",
         font_size=11, color=SLATE_500, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: Claims Data Requirements
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 4, Inches(0.6), Inches(0.4))
add_category(slide, "Data Schema", Inches(1.15), Inches(0.45))
add_question_title(slide, "What data fields does Pear Suite require to submit\na Medi-Cal claim for a CHW encounter?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.2),
                "Why This Matters", [
                    "Compass must capture every field",
                    "during the CHW session.",
                    "Missing fields = rejected claims",
                    "= lost revenue for the org."
                ])

add_highlight_box(slide, Inches(0.6), Inches(3.7), Inches(3.5), Inches(1.0),
                  "Our Goal",
                  "Map Pear Suite's required schema to Compass\nform fields → zero manual re-entry.")

# 4 data boxes
data_boxes = [
    ("Member Information", GREEN_BG, COMPASS_GREEN, COMPASS_DARK,
     [("Medi-Cal ID", True), ("Member name & DOB", True), ("Address / County", True),
      ("Managed care plan", False), ("Eligibility status", False)]),
    ("Service Details", GOLD_BG, PEAR_GOLD, PEAR_DARK,
     [("CPT / HCPCS codes", False), ("ICD-10 diagnosis codes", False),
      ("Service date & duration", True), ("Place of service", True), ("Rendering provider NPI", False)]),
    ("Provider / Org", BLUE_BG, BLUE_500, BLUE_DARK,
     [("Billing provider NPI", False), ("Tax ID / EIN", False), ("Org name & address", True)]),
    ("Session Documentation", PURPLE_BG, PURPLE_500, PURPLE_DARK,
     [("Session notes / summary", True), ("Goals addressed", True), ("Referrals made", True)]),
]

positions = [
    (Inches(4.6), Inches(2.2), Inches(3.8), Inches(2.4)),
    (Inches(8.8), Inches(2.2), Inches(3.8), Inches(2.4)),
    (Inches(4.6), Inches(4.8), Inches(3.8), Inches(1.8)),
    (Inches(8.8), Inches(4.8), Inches(3.8), Inches(1.8)),
]

for (title, bg, bc, tc, fields), (bx, by, bw, bh) in zip(data_boxes, positions):
    add_shape(slide, bx, by, bw, bh, fill=bg, border_color=bc)
    add_text(slide, bx + Inches(0.15), by + Inches(0.1), bw - Inches(0.3), Inches(0.25),
             title, font_size=13, bold=True, color=tc)
    for fi, (fname, captured) in enumerate(fields):
        fy = by + Inches(0.45) + fi * Inches(0.28)
        add_text(slide, bx + Inches(0.15), fy, Inches(2.8), Inches(0.22),
                 fname, font_size=11, color=SLATE_700)
        dot_color = COMPASS_GREEN if captured else PEAR_GOLD
        dot_label = "C" if captured else "?"
        add_circle(slide, bx + bw - Inches(0.45), fy + Inches(0.02), Inches(0.2),
                   dot_color, dot_label, 8, WHITE)

# Legend
add_circle(slide, Inches(0.6), Inches(6.8), Inches(0.2), COMPASS_GREEN, "C", 8, WHITE)
add_text(slide, Inches(0.85), Inches(6.78), Inches(2.5), Inches(0.22),
         "= Compass captures today", font_size=10, color=SLATE_700)
add_circle(slide, Inches(3.0), Inches(6.8), Inches(0.2), PEAR_GOLD, "?", 8, WHITE)
add_text(slide, Inches(3.25), Inches(6.78), Inches(2.5), Inches(0.22),
         "= Need to confirm with Pear Suite", font_size=10, color=SLATE_700)


# ============================================================
# SLIDE 6: Authentication
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 5, Inches(0.6), Inches(0.4))
add_category(slide, "Security", Inches(1.15), Inches(0.45))
add_question_title(slide, "How do you handle authentication for third-party integrations\n— API keys, OAuth, or org-scoped tokens?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.2),
                "Why This Matters", [
                    "Auth model determines integration",
                    "architecture. Per-org tokens vs.",
                    "per-user OAuth = very different",
                    "implementation patterns."
                ])

add_context_box(slide, Inches(0.6), Inches(3.7), Inches(3.5), Inches(1.3),
                "Compass Auth Stack", [
                    "AWS Cognito (user auth)",
                    "JWT tokens",
                    "Role-based access control",
                    "Per-org data isolation"
                ])

# 3 auth option rows
auth_opts = [
    ("A: API Key (per-org)", "Simplest. Key stored in AWS Secrets Manager.",
     BLUE_BG, BLUE_500, BLUE_DARK),
    ("B: OAuth 2.0", "Per-user. More complex. Standard flow.",
     PURPLE_BG, PURPLE_500, PURPLE_DARK),
    ("C: Service-to-Service (mTLS / JWT)", "Most secure. Needs cert management.",
     GREEN_BG, COMPASS_GREEN, COMPASS_DARK),
]
for i, (title, desc, bg, bc, tc) in enumerate(auth_opts):
    y = Inches(2.3) + i * Inches(1.35)
    add_shape(slide, Inches(4.8), y, Inches(7.8), Inches(1.1), fill=bg, border_color=bc)
    add_text(slide, Inches(5.0), y + Inches(0.12), Inches(5), Inches(0.3),
             title, font_size=14, bold=True, color=tc)

    # Mini flow: Compass → [method] → Pear Suite
    flow_y2 = y + Inches(0.55)
    add_shape(slide, Inches(5.0), flow_y2, Inches(1.2), Inches(0.35),
              fill=WHITE, border_color=COMPASS_GREEN)
    add_text(slide, Inches(5.0), flow_y2 + Inches(0.02), Inches(1.2), Inches(0.3),
             "Compass", font_size=10, color=SLATE_700, alignment=PP_ALIGN.CENTER)
    add_right_arrow_shape(slide, Inches(6.3), flow_y2 + Inches(0.05), Inches(0.3), Inches(0.22), bc)
    add_shape(slide, Inches(6.7), flow_y2, Inches(1.2), Inches(0.35),
              fill=WHITE, border_color=PEAR_GOLD)
    add_text(slide, Inches(6.7), flow_y2 + Inches(0.02), Inches(1.2), Inches(0.3),
             "Pear Suite", font_size=10, color=SLATE_700, alignment=PP_ALIGN.CENTER)

    add_text(slide, Inches(8.5), flow_y2 - Inches(0.05), Inches(3.5), Inches(0.4),
             desc, font_size=11, color=SLATE_700)

# Question pill
add_shape(slide, Inches(6.5), Inches(6.3), Inches(3.5), Inches(0.45),
          fill=SLATE_900, border_color=None)
add_text(slide, Inches(6.5), Inches(6.33), Inches(3.5), Inches(0.4),
         "Which model does Pear Suite use?", font_size=13, bold=True, color=WHITE,
         alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: HIPAA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 6, Inches(0.6), Inches(0.4))
add_category(slide, "Compliance", Inches(1.15), Inches(0.45))
add_question_title(slide, "What are your HIPAA requirements for integration partners\n— do you require a BAA and security review?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.2),
                "Why This Matters", [
                    "PHI will flow between both systems.",
                    "Need Pear Suite's compliance bar",
                    "to build a concrete checklist.",
                ])

add_highlight_box(slide, Inches(0.6), Inches(3.7), Inches(3.5), Inches(1.0),
                  "Our Current State",
                  "Compass audit flagged HIPAA at 1/10.\nKnowing Pear Suite's bar gives us\na concrete compliance checklist.",
                  bg=GOLD_BG, border=PEAR_GOLD, label_color=PEAR_DARK, text_color=PEAR_DARK)

# Compass compliance box
add_shape(slide, Inches(4.8), Inches(2.2), Inches(3.2), Inches(4.5),
          fill=GREEN_BG, border_color=COMPASS_GREEN, border_width=Pt(2))
add_text(slide, Inches(4.8), Inches(2.3), Inches(3.2), Inches(0.3),
         "Compass", font_size=14, bold=True, color=COMPASS_DARK, alignment=PP_ALIGN.CENTER)

compass_items = [
    ("AWS (BAA needed)", PEAR_GOLD), ("Encryption at rest", COMPASS_GREEN),
    ("Encryption in transit", COMPASS_GREEN), ("Audit logging", RED_500),
    ("Access controls", COMPASS_GREEN), ("Data retention", RED_500),
]
for i, (item, status) in enumerate(compass_items):
    y = Inches(2.8) + i * Inches(0.4)
    add_text(slide, Inches(5.0), y, Inches(2.2), Inches(0.25), item, font_size=11, color=SLATE_700)
    add_circle(slide, Inches(7.5), y + Inches(0.05), Inches(0.18), status, "", 1, status)

# PHI Exchange box
add_shape(slide, Inches(8.2), Inches(3.5), Inches(1.2), Inches(1.2),
          fill=SLATE_900, border_color=None)
add_text(slide, Inches(8.2), Inches(3.7), Inches(1.2), Inches(0.3),
         "PHI", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.2), Inches(4.0), Inches(1.2), Inches(0.3),
         "Exchange", font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.2), Inches(4.3), Inches(1.2), Inches(0.25),
         "TLS 1.3", font_size=9, color=SLATE_300, alignment=PP_ALIGN.CENTER)

# Pear Suite compliance box
add_shape(slide, Inches(9.6), Inches(2.2), Inches(3.2), Inches(4.5),
          fill=GOLD_BG, border_color=PEAR_GOLD, border_width=Pt(2))
add_text(slide, Inches(9.6), Inches(2.3), Inches(3.2), Inches(0.3),
         "Pear Suite", font_size=14, bold=True, color=PEAR_DARK, alignment=PP_ALIGN.CENTER)

ps_items = ["BAA required?", "Security review?", "SOC 2?",
            "Pen test needed?", "Data handling SLA?", "Breach protocol?"]
for i, item in enumerate(ps_items):
    y = Inches(2.8) + i * Inches(0.4)
    add_text(slide, Inches(9.8), y, Inches(2.2), Inches(0.25), item, font_size=11, color=SLATE_700)
    add_circle(slide, Inches(12.3), y + Inches(0.05), Inches(0.18), SLATE_500, "", 1, SLATE_500)

# Legend
legend_items = [
    (COMPASS_GREEN, "Done"), (PEAR_GOLD, "In progress"),
    (RED_500, "To do"), (SLATE_500, "Ask Nick"),
]
for i, (c, label) in enumerate(legend_items):
    x = Inches(4.8) + i * Inches(1.8)
    add_circle(slide, x, Inches(7.0), Inches(0.18), c, "", 1, c)
    add_text(slide, x + Inches(0.25), Inches(6.97), Inches(1.2), Inches(0.22),
             label, font_size=10, color=SLATE_700)


# ============================================================
# SLIDE 8: Claims Lifecycle
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 7, Inches(0.6), Inches(0.4))
add_category(slide, "Workflow", Inches(1.15), Inches(0.45))
add_question_title(slide, "What does the claims lifecycle look like after submission\n— how would Compass get status updates?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.2),
                "Why This Matters", [
                    "If Compass submits but never knows",
                    "if the claim paid out, CHW orgs",
                    "lose visibility. Closing the loop",
                    "is the entire value prop."
                ])

add_highlight_box(slide, Inches(0.6), Inches(3.7), Inches(3.5), Inches(1.0),
                  "Ideal State",
                  "CHW sees claim status in-app:\nsubmitted → accepted → denied → paid.\nZero context-switching.")

# Lifecycle steps (vertical timeline)
lifecycle = [
    ("1", "Submitted", "Compass sends session data to Pear Suite via API", COMPASS_GREEN, GREEN_BG),
    ("2", "Validated / Scrubbed", "Pear Suite validates claim against Medi-Cal rules", BLUE_500, BLUE_BG),
    ("3", "Clearinghouse Processing", "Claim routed to Medi-Cal / managed care plan", PURPLE_500, PURPLE_BG),
    ("4", "Adjudicated — Paid / Denied / Resubmit", "How does Compass learn the outcome? Webhook? Polling?", PEAR_GOLD, GOLD_BG),
]
for i, (num, title, desc, color, bg) in enumerate(lifecycle):
    y = Inches(2.2) + i * Inches(1.2)
    # Circle
    add_circle(slide, Inches(4.8), y + Inches(0.08), Inches(0.4), color, num, 13, WHITE)
    # Box
    add_shape(slide, Inches(5.4), y, Inches(7.0), Inches(0.85), fill=bg, border_color=color)
    add_text(slide, Inches(5.6), y + Inches(0.1), Inches(6.5), Inches(0.3),
             title, font_size=14, bold=True, color=SLATE_900)
    add_text(slide, Inches(5.6), y + Inches(0.45), Inches(6.5), Inches(0.25),
             desc, font_size=11, color=SLATE_700)
    # Vertical line between
    if i < 3:
        line_y = y + Inches(0.85)
        add_shape(slide, Inches(4.95), line_y, Inches(0.08), Inches(0.35),
                  fill=SLATE_300, border_color=None, shape_type=MSO_SHAPE.RECTANGLE)

# Callback arrow label
add_text(slide, Inches(10.0), Inches(6.5), Inches(2.5), Inches(0.25),
         "↩ Status callback?", font_size=12, bold=True, color=RED_500)


# ============================================================
# SLIDE 9: Sandbox
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 8, Inches(0.6), Inches(0.4))
add_category(slide, "Development", Inches(1.15), Inches(0.45))
add_question_title(slide, "Is there a sandbox or staging environment\nwe could develop and test against?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.3),
                "Why This Matters", [
                    "We need a non-production environment",
                    "to build, test, and iterate.",
                    "No sandbox = significantly longer",
                    "development timeline."
                ])

add_context_box(slide, Inches(0.6), Inches(3.8), Inches(3.5), Inches(1.5),
                "What We Need", [
                    "API endpoint(s) with test data",
                    "Test Medi-Cal member IDs",
                    "Mock claim responses",
                    "Rate limits / quotas",
                    "API documentation"
                ])

# 3 environment boxes
envs = [
    ("Development", "local → sandbox\nTest data / mock claims", PURPLE_BG, PURPLE_500, PURPLE_DARK),
    ("Staging", "staging → staging\nEnd-to-end validation", BLUE_BG, BLUE_500, BLUE_DARK),
    ("Production", "prod → prod\nReal PHI / real claims", GREEN_BG, COMPASS_GREEN, COMPASS_DARK),
]
# Top row: dev and staging
for i, (title, desc, bg, bc, tc) in enumerate(envs[:2]):
    x = Inches(5.0) + i * Inches(3.8)
    add_shape(slide, x, Inches(2.3), Inches(3.3), Inches(2.0),
              fill=bg, border_color=bc, border_width=Pt(1.5))
    add_text(slide, x, Inches(2.4), Inches(3.3), Inches(0.3),
             title, font_size=14, bold=True, color=tc, alignment=PP_ALIGN.CENTER)
    # Mini compass → pear suite
    add_shape(slide, x + Inches(0.3), Inches(2.9), Inches(1.0), Inches(0.4),
              fill=WHITE, border_color=COMPASS_GREEN)
    add_text(slide, x + Inches(0.3), Inches(2.92), Inches(1.0), Inches(0.35),
             "Compass", font_size=10, color=SLATE_700, alignment=PP_ALIGN.CENTER)
    add_right_arrow_shape(slide, x + Inches(1.35), Inches(3.0), Inches(0.25), Inches(0.2), bc)
    add_shape(slide, x + Inches(1.7), Inches(2.9), Inches(1.0), Inches(0.4),
              fill=WHITE, border_color=PEAR_GOLD)
    add_text(slide, x + Inches(1.7), Inches(2.92), Inches(1.0), Inches(0.35),
             "Pear Suite", font_size=10, color=SLATE_700, alignment=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(3.6), Inches(3.3), Inches(0.5),
             desc, font_size=10, color=tc, alignment=PP_ALIGN.CENTER)

# Production (bottom, centered)
prod = envs[2]
add_shape(slide, Inches(6.9), Inches(4.8), Inches(3.3), Inches(2.2),
          fill=prod[2], border_color=prod[3], border_width=Pt(2.5))
add_text(slide, Inches(6.9), Inches(4.9), Inches(3.3), Inches(0.3),
         prod[0], font_size=14, bold=True, color=prod[4], alignment=PP_ALIGN.CENTER)
add_shape(slide, Inches(7.2), Inches(5.4), Inches(1.0), Inches(0.4),
          fill=WHITE, border_color=COMPASS_GREEN, border_width=Pt(1.5))
add_text(slide, Inches(7.2), Inches(5.42), Inches(1.0), Inches(0.35),
         "Compass", font_size=10, color=SLATE_700, alignment=PP_ALIGN.CENTER)
add_right_arrow_shape(slide, Inches(8.25), Inches(5.5), Inches(0.25), Inches(0.2), COMPASS_GREEN)
add_shape(slide, Inches(8.6), Inches(5.4), Inches(1.0), Inches(0.4),
          fill=WHITE, border_color=PEAR_GOLD, border_width=Pt(1.5))
add_text(slide, Inches(8.6), Inches(5.42), Inches(1.0), Inches(0.35),
         "Pear Suite", font_size=10, color=SLATE_700, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(6.9), Inches(6.2), Inches(3.3), Inches(0.5),
         "Real PHI / real claims", font_size=10, color=COMPASS_DARK, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 10: Reference Integrations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 9, Inches(0.6), Inches(0.4))
add_category(slide, "Precedent", Inches(1.15), Inches(0.45))
add_question_title(slide, "Are there other platforms integrating with Pear Suite today\nthat we could reference as a model?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.2),
                "Why This Matters", [
                    "If someone has already done this,",
                    "we can learn from their patterns.",
                    "Saves months of guesswork."
                ])

add_highlight_box(slide, Inches(0.6), Inches(3.7), Inches(3.5), Inches(1.0),
                  "What We're Looking For",
                  "Integration patterns, API usage examples,\ncommon pitfalls, time-to-integration\nbenchmarks from existing partners.")

# Hub-and-spoke diagram
# Center: Pear Suite
center_x, center_y = Inches(8.5), Inches(4.0)
add_circle(slide, center_x - Inches(0.6), center_y - Inches(0.6), Inches(1.2),
           PEAR_GOLD, "Pear\nSuite", 13, WHITE)

# Compass (top)
add_circle(slide, center_x - Inches(0.45), center_y - Inches(2.4), Inches(0.9),
           COMPASS_GREEN, "Compass\nCHW", 10, WHITE)

# Partner placeholders
partners = [
    ("EHR\nPartner?", BLUE_500, -2.2, -0.6),
    ("Health\nPlan?", PURPLE_500, 2.2, -0.6),
    ("Care\nMgmt?", RED_500, -2.0, 1.4),
    ("Other\nCHW?", COMPASS_GREEN, 2.0, 1.4),
]
for label, color, dx, dy in partners:
    px = center_x + Inches(dx) - Inches(0.45)
    py = center_y + Inches(dy) - Inches(0.45)
    add_circle(slide, px, py, Inches(0.9), color, label, 9, WHITE)

add_text(slide, Inches(5.5), Inches(6.6), Inches(6), Inches(0.3),
         "Dashed partners are unknown — ask Nick to fill in the ecosystem",
         font_size=11, bold=True, color=SLATE_500, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 11: Phase 1 Roadmap
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_white_slide(slide)
add_q_badge(slide, 10, Inches(0.6), Inches(0.4))
add_category(slide, "Next Steps", Inches(1.15), Inches(0.45))
add_question_title(slide, "What would a Phase 1 integration look like,\nand what's the fastest path to a working demo?",
                   Inches(0.6), Inches(1.0))

add_context_box(slide, Inches(0.6), Inches(2.3), Inches(3.5), Inches(1.0),
                "Why This Matters", [
                    "Closes with a concrete next step.",
                    "Ship something small and real,",
                    "not plan forever."
                ])

add_highlight_box(slide, Inches(0.6), Inches(3.5), Inches(3.5), Inches(0.9),
                  "Our Proposal",
                  "One CHW org, one session type, one claim\nflow — end to end. Prove it, then scale.")

# Timeline with 3 phases
# Horizontal line
add_shape(slide, Inches(4.8), Inches(2.15), Inches(7.8), Inches(0.04),
          fill=SLATE_300, border_color=None, shape_type=MSO_SHAPE.RECTANGLE)

# Phase circles on timeline
phases_pos = [Inches(6.2), Inches(8.8), Inches(11.4)]
phase_colors = [COMPASS_GREEN, BLUE_500, PURPLE_500]
phase_labels = ["1", "2", "3"]
phase_titles = ["Phase 1", "Phase 2", "Phase 3"]
for i in range(3):
    add_circle(slide, phases_pos[i] - Inches(0.2), Inches(1.95), Inches(0.45),
               phase_colors[i], phase_labels[i], 14, WHITE)
    add_text(slide, phases_pos[i] - Inches(0.5), Inches(1.65), Inches(1), Inches(0.25),
             phase_titles[i], font_size=11, bold=True, color=phase_colors[i],
             alignment=PP_ALIGN.CENTER)

# Phase 1 box
add_shape(slide, Inches(4.8), Inches(2.6), Inches(3.2), Inches(4.5),
          fill=GREEN_BG, border_color=COMPASS_GREEN, border_width=Pt(2.5))
add_text(slide, Inches(4.8), Inches(2.7), Inches(3.2), Inches(0.3),
         "MVP Integration", font_size=14, bold=True, color=COMPASS_DARK, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(4.8), Inches(3.0), Inches(3.2), Inches(0.25),
         "4–6 weeks", font_size=11, color=SLATE_500, alignment=PP_ALIGN.CENTER)

p1_items = ["1 CHW organization", "1 session type", "Submit claims via API",
            "Basic status callback", "BAA signed"]
for i, item in enumerate(p1_items):
    add_text(slide, Inches(5.1), Inches(3.5) + i * Inches(0.32), Inches(2.5), Inches(0.25),
             f"• {item}", font_size=11, color=SLATE_700)

add_text(slide, Inches(5.1), Inches(5.2), Inches(2.5), Inches(0.25),
         "Deliverable:", font_size=11, bold=True, color=COMPASS_DARK)
add_text(slide, Inches(5.1), Inches(5.5), Inches(2.5), Inches(0.5),
         "Working demo with\nreal claim submission", font_size=11, color=SLATE_700)

# Phase 2 box
add_shape(slide, Inches(8.2), Inches(2.6), Inches(2.6), Inches(4.0),
          fill=BLUE_BG, border_color=BLUE_500)
add_text(slide, Inches(8.2), Inches(2.7), Inches(2.6), Inches(0.3),
         "Expand", font_size=14, bold=True, color=BLUE_DARK, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(8.2), Inches(3.0), Inches(2.6), Inches(0.25),
         "8–12 weeks", font_size=11, color=SLATE_500, alignment=PP_ALIGN.CENTER)

p2_items = ["Multi-org support", "All session types", "Full status tracking",
            "Denial management", "Dashboard analytics"]
for i, item in enumerate(p2_items):
    add_text(slide, Inches(8.4), Inches(3.5) + i * Inches(0.32), Inches(2.2), Inches(0.25),
             f"• {item}", font_size=11, color=SLATE_700)

add_text(slide, Inches(8.4), Inches(5.2), Inches(2.2), Inches(0.25),
         "Deliverable:", font_size=11, bold=True, color=BLUE_DARK)
add_text(slide, Inches(8.4), Inches(5.5), Inches(2.2), Inches(0.25),
         "Production-ready", font_size=11, color=SLATE_700)

# Phase 3 box
add_shape(slide, Inches(11.0), Inches(2.6), Inches(1.8), Inches(3.5),
          fill=PURPLE_BG, border_color=PURPLE_500)
add_text(slide, Inches(11.0), Inches(2.7), Inches(1.8), Inches(0.3),
         "Scale", font_size=14, bold=True, color=PURPLE_DARK, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(11.0), Inches(3.0), Inches(1.8), Inches(0.25),
         "Ongoing", font_size=11, color=SLATE_500, alignment=PP_ALIGN.CENTER)

p3_items = ["EHR sync", "Auto-coding", "Reporting", "Multi-state"]
for i, item in enumerate(p3_items):
    add_text(slide, Inches(11.2), Inches(3.5) + i * Inches(0.32), Inches(1.4), Inches(0.25),
             f"• {item}", font_size=11, color=SLATE_700)


# ============================================================
# SLIDE 12: Closing
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
make_dark_slide(slide)

add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
         "Let's Build This Together", font_size=44, bold=True, color=WHITE,
         alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(3.5), Inches(8), Inches(1.2),
         "Compass handles the CHW experience.\nPear Suite handles billing infrastructure.\nTogether: field visit to paid claim, seamlessly.",
         font_size=20, color=SLATE_300, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(5.5), Inches(8), Inches(0.4),
         "Next steps: sandbox access  •  BAA discussion  •  schema alignment",
         font_size=14, color=SLATE_500, alignment=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────
output_path = "/Users/akrammahmoud/Desktop/Projects/Compass/docs/Compass_x_PearSuite_Meeting.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
